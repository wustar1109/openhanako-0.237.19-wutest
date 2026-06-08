#!/usr/bin/env python3
from __future__ import annotations

import argparse
import shutil
import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from office_common import (
    DocumentSkillError,
    json_error,
    json_success,
    load_ops,
    local_name,
    numeric_part_key,
    prepare_output,
    replace_text_in_containers,
    require_file,
    write_modified_zip,
    write_text_or_json,
    xml_bytes,
)


def apply_replace_text(input_path: Path, output_path: Path, op: dict) -> dict:
    find = op.get("find")
    replace = op.get("replace")
    if not isinstance(find, str) or not isinstance(replace, str):
        raise DocumentSkillError("replace_text requires string find and replace fields.", code="invalid_operation")
    max_count = op.get("count")
    if max_count is not None:
        max_count = int(max_count)
        if max_count < 1:
            raise DocumentSkillError("replace_text count must be positive.", code="invalid_operation")

    replacements: dict[str, bytes] = {}
    total = 0
    with zipfile.ZipFile(input_path) as zf:
        parts = sorted(
            [
                name for name in zf.namelist()
                if (name.startswith("ppt/slides/slide") or name.startswith("ppt/notesSlides/notesSlide"))
                and name.endswith(".xml")
            ],
            key=numeric_part_key,
        )
        for name in parts:
            try:
                root = ET.fromstring(zf.read(name))
            except Exception:
                continue
            changed = replace_text_in_containers(
                root,
                container_names={"p"},
                find=find,
                replace=replace,
                max_count=max_count if total == 0 else (None if max_count is None else max_count - total),
            )
            if changed:
                replacements[name] = xml_bytes(root)
                total += changed
                if max_count is not None and total >= max_count:
                    break

    if not replacements:
        shutil.copyfile(input_path, output_path)
    else:
        write_modified_zip(input_path, output_path, replacements)
    return {"op": "replace_text", "matches": total}


def require_python_pptx():
    try:
        import pptx  # type: ignore
        return pptx
    except Exception as exc:
        raise DocumentSkillError(
            "This PPTX operation requires the optional dependency python-pptx.",
            code="missing_dependency",
            details={"package": "python-pptx", "original_error": str(exc)},
        ) from exc


def apply_python_pptx_op(input_path: Path, output_path: Path, op: dict) -> dict:
    pptx = require_python_pptx()
    from pptx.util import Inches  # type: ignore

    prs = pptx.Presentation(str(input_path))
    kind = op["op"]
    slide_number = int(op.get("slide", 0))
    if slide_number < 1 or slide_number > len(prs.slides):
        raise DocumentSkillError(f"Slide {slide_number} is out of range.", code="slide_out_of_range")
    slide = prs.slides[slide_number - 1]

    if kind == "set_shape_text":
        text = op.get("text")
        if not isinstance(text, str):
            raise DocumentSkillError("set_shape_text requires text.", code="invalid_operation")
        shape = None
        if "shape_index" in op:
            idx = int(op["shape_index"])
            if idx < 1 or idx > len(slide.shapes):
                raise DocumentSkillError(f"shape_index {idx} is out of range.", code="shape_out_of_range")
            shape = slide.shapes[idx - 1]
        elif "shape_name" in op:
            target = str(op["shape_name"])
            shape = next((s for s in slide.shapes if getattr(s, "name", "") == target), None)
            if shape is None:
                raise DocumentSkillError(f"Shape not found by name: {target}", code="shape_not_found")
        else:
            raise DocumentSkillError("set_shape_text requires shape_index or shape_name.", code="invalid_operation")
        if not hasattr(shape, "text"):
            raise DocumentSkillError("Target shape does not expose text.", code="not_text_shape")
        shape.text = text
        result = {"op": kind, "slide": slide_number}
    elif kind == "add_textbox":
        text = op.get("text")
        if not isinstance(text, str):
            raise DocumentSkillError("add_textbox requires text.", code="invalid_operation")
        left = Inches(float(op.get("left", 1)))
        top = Inches(float(op.get("top", 1)))
        width = Inches(float(op.get("width", 8)))
        height = Inches(float(op.get("height", 1)))
        box = slide.shapes.add_textbox(left, top, width, height)
        box.text = text
        result = {"op": kind, "slide": slide_number}
    else:
        raise DocumentSkillError(f"Unsupported PPTX operation: {kind}", code="unsupported_operation")

    prs.save(str(output_path))
    return result


def run_ops(input_path: Path, output_path: Path, ops: list[dict]) -> list[dict]:
    tmp_dir = Path(tempfile.mkdtemp(prefix="edit-pptx-"))
    results: list[dict] = []
    try:
        current = tmp_dir / "step-0.pptx"
        shutil.copyfile(input_path, current)
        for idx, op in enumerate(ops, 1):
            next_path = tmp_dir / f"step-{idx}.pptx"
            kind = op["op"]
            if kind == "replace_text":
                results.append(apply_replace_text(current, next_path, op))
            elif kind in {"set_shape_text", "add_textbox"}:
                results.append(apply_python_pptx_op(current, next_path, op))
            else:
                raise DocumentSkillError(f"Unsupported PPTX operation: {kind}", code="unsupported_operation")
            current = next_path
        shutil.copyfile(current, output_path)
        return results
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def main() -> int:
    parser = argparse.ArgumentParser(description="Apply safe, limited edits to PPTX files.")
    parser.add_argument("input")
    parser.add_argument("output")
    parser.add_argument("--ops", required=True, help="JSON operations file.")
    args = parser.parse_args()

    try:
        input_path = require_file(args.input)
        output_path = prepare_output(input_path, args.output)
        if input_path.suffix.lower() != ".pptx" or output_path.suffix.lower() != ".pptx":
            raise DocumentSkillError("edit_pptx.py requires .pptx input and output paths.", code="unsupported_format")
        ops = load_ops(args.ops)
        results = run_ops(input_path, output_path, ops)
        write_text_or_json(None, json_success(input=str(input_path), output=str(output_path), operations=results))
        return 0
    except Exception as exc:
        write_text_or_json(None, json_error(exc))
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
